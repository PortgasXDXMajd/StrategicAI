import os
import pdfplumber
import pandas as pd
from io import BytesIO
from docx import Document
from pptx import Presentation
from fastapi import UploadFile


class FileHelper:
    @staticmethod
    def extract_from_txt(file: BytesIO) -> str:
        return file.read().decode("utf-8")

    @staticmethod
    def extract_from_csv(file: BytesIO) -> str:
        df = pd.read_csv(file)
        return "\n".join([" ".join(map(str, row)) for row in df.values.tolist()])

    @staticmethod
    def extract_from_excel(file: BytesIO) -> str:
        df = pd.read_excel(file)
        return "\n".join([" ".join(map(str, row)) for row in df.values.tolist()])

    @staticmethod
    def extract_from_pdf(file: BytesIO) -> str:
        text = ""
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                text += (page.extract_text() or "") + "\n"
        return text.strip()

    @staticmethod
    def extract_from_word(file: BytesIO) -> str:
        doc = Document(file)
        return "\n".join([para.text for para in doc.paragraphs if para.text]).strip()

    @staticmethod
    def extract_from_ppt(file: BytesIO) -> str:
        presentation = Presentation(file)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()

    @staticmethod
    def extract_full_content(file: UploadFile) -> str:
        file_extension = os.path.splitext(file.filename)[1].lower()
        # Read file into BytesIO for in-memory handling
        file_data = BytesIO(file.file.read())
        if file_extension == ".txt":
            return FileHelper.extract_from_txt(file_data)
        elif file_extension == ".csv":
            return FileHelper.extract_from_csv(file_data)
        elif file_extension in [".xls", ".xlsx"]:
            return FileHelper.extract_from_excel(file_data)
        elif file_extension == ".pdf":
            return FileHelper.extract_from_pdf(file_data)
        elif file_extension == ".docx":
            return FileHelper.extract_from_word(file_data)
        elif file_extension == ".pptx":
            return FileHelper.extract_from_ppt(file_data)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
